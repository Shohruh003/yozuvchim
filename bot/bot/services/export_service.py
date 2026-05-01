import re
import asyncio
from datetime import datetime
from pathlib import Path

import docx
from docx.shared import Pt as DocxPt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING, WD_TAB_LEADER
from docx.oxml import OxmlElement
from docx.oxml.ns import qn

from pptx import Presentation
from pptx.util import Pt as PptxPt

from ..config import logger


class ExportService:
    # ---------------------------
    # DOCX
    # ---------------------------
    @staticmethod
    async def to_docx(title: str, content: str, path: Path, university_header: str = None, doc_type: str = "coursework", meta: dict = None) -> bool:
        return await asyncio.to_thread(ExportService._to_docx_sync, title, content, path, university_header, doc_type, meta)

    @staticmethod
    def _to_docx_sync(title: str, content: str, path: Path, university_header: str = None, doc_type: str = "coursework", meta: dict = None) -> bool:
        try:
            doc = docx.Document()

            # Margins (O'AK-ish)
            for section in doc.sections:
                section.top_margin = Cm(2.0)
                section.bottom_margin = Cm(2.0)
                section.left_margin = Cm(3.0)
                section.right_margin = Cm(1.5)

            # Normal style (body)
            normal = doc.styles["Normal"]
            normal.font.name = "Times New Roman"
            normal._element.rPr.rFonts.set(qn("w:eastAsia"), "Times New Roman")
            normal.font.size = DocxPt(14)

            pf = normal.paragraph_format
            pf.line_spacing_rule = WD_LINE_SPACING.ONE_POINT_FIVE
            pf.space_after = DocxPt(0)
            pf.first_line_indent = Cm(1.25)
            pf.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

            # Make headings O'AK-compliant: 14pt bold, no indent
            for lvl in ("Heading 1", "Heading 2", "Heading 3", "Title"):
                if lvl in doc.styles:
                    st = doc.styles[lvl]
                    st.font.name = "Times New Roman"
                    st._element.rPr.rFonts.set(qn("w:eastAsia"), "Times New Roman")
                    st.font.size = DocxPt(14)
                    st.font.bold = True
                    st.font.color.rgb = RGBColor(0, 0, 0)
                    st.paragraph_format.first_line_indent = Cm(0)
                    st.paragraph_format.space_before = DocxPt(12)
                    st.paragraph_format.space_after = DocxPt(6)
                    st.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER

            meta = meta or {}
            clean_title = title.strip()

            # Title page — tezis and article get inline header, others get full title page
            if doc_type == "thesis":
                # Tezis: compact inline header (no title page)
                ExportService._render_thesis_header(doc, clean_title, meta)
            elif doc_type == "article":
                # Article: author info (right-aligned) + title
                ExportService._render_article_header(doc, clean_title, meta)
            elif meta.get("uni") or meta.get("subject"):
                ExportService._render_title_page(doc, clean_title, meta, doc_type)
                doc.add_page_break()
            elif university_header:
                h = doc.add_paragraph(university_header)
                h.alignment = WD_ALIGN_PARAGRAPH.CENTER
                doc.add_page_break()
            else:
                t = doc.add_paragraph(ExportService._clean_text(clean_title))
                t.style = doc.styles["Title"]
                t.alignment = WD_ALIGN_PARAGRAPH.CENTER

            # Page numbers (footer)
            ExportService._add_page_numbers(doc)

            # Pre-process: extract markdown tables from content before block splitting
            # Replace each table with a placeholder <<<TABLE_N>>>
            _table_store = {}
            _table_counter = [0]

            def _extract_tables(text):
                """Find all markdown table regions and replace with placeholders."""
                lines = text.split("\n")
                result = []
                i = 0
                while i < len(lines):
                    ln = lines[i].strip()
                    # A table starts with a line that has multiple | and next line is separator
                    if "|" in ln and i + 1 < len(lines):
                        next_ln = lines[i + 1].strip()
                        is_sep = bool(re.match(r"^[ \t]*\|?[\-\|\s:]+\|?[ \t]*$", next_ln)) and "-" in next_ln
                        if is_sep:
                            # Collect all lines of this table
                            table_lines = []
                            while i < len(lines) and ("|" in lines[i] or lines[i].strip() == ""):
                                if lines[i].strip():
                                    table_lines.append(lines[i])
                                i += 1
                            key = f"<<<TABLE_{_table_counter[0]}>>>"
                            _table_store[key] = "\n".join(table_lines)
                            _table_counter[0] += 1
                            result.append("")  # blank line before
                            result.append(key)
                            result.append("")  # blank line after
                            continue
                    result.append(lines[i])
                    i += 1
                return "\n".join(result)

            processed_content = _extract_tables(content or "")
            blocks = [b.strip() for b in processed_content.split("\n\n") if b.strip()]

            # TOC Section - Static list for better compatibility
            if doc_type not in ("taqdimot", "article", "thesis") and len(blocks) > 3:
                lang = (meta.get("language") or "uz").lower()
                toc_title = {"en": "TABLE OF CONTENTS", "ru": "СОДЕРЖАНИЕ"}.get(lang, "MUNDARIJA")
                h = doc.add_paragraph(toc_title)
                h.style = doc.styles["Heading 1"]
                h.alignment = WD_ALIGN_PARAGRAPH.CENTER
                h.paragraph_format.space_after = DocxPt(12)
                
                # Extract headings for static TOC
                for blk in blocks:
                    if blk.startswith("# ") or blk.startswith("## ") or blk.startswith("### "):
                        level = 0
                        if blk.startswith("### "): level = 2
                        elif blk.startswith("## "): level = 1
                        
                        title_text = blk.lstrip("#").strip()
                        # Clean title_text for TOC (removing any bold/italic artifacts)
                        title_text = title_text.replace("**", "").replace("__", "").replace("*", "").replace("_", "").replace("`", "")
                        
                        p = doc.add_paragraph()
                        p.alignment = WD_ALIGN_PARAGRAPH.LEFT
                        p.paragraph_format.first_line_indent = Cm(0.5 * level)
                        p.paragraph_format.space_after = DocxPt(0)

                        # Add tab stops for leader dots (TOC style)
                        tabs = p.paragraph_format.tab_stops
                        tabs.add_tab_stop(Cm(16), alignment=WD_ALIGN_PARAGRAPH.RIGHT, leader=WD_TAB_LEADER.DOTS)

                        run = p.add_run(title_text)
                        run.font.size = DocxPt(12)

                        # Add a tab character to trigger leader dots
                        p.add_run("\t")
                
                doc.add_page_break()

            in_references = False
            refs_started = False  # For combined sections (keywords+refs)
            ref_counter = 0
            prev_heading = ""

            def _check_refs_heading(heading_text):
                """Check if heading is a references section. Returns (in_refs, is_combined)."""
                hc = re.sub(r'^[\d]+[\.\)]\s*', '', heading_text).strip().lower()
                # Pure references headings
                if ("foydalanilgan" in hc and "adabiyot" in hc) or hc.startswith("references") or \
                   ("список" in hc and "литератур" in hc) or hc.startswith("список литератур"):
                    return True, False
                # Combined: keywords + references (thesis)
                if ("kalit" in hc and "adabiyot" in hc) or \
                   ("keyword" in hc and "reference" in hc) or \
                   ("ключев" in hc and "литератур" in hc):
                    return True, True
                # Standalone "Adabiyotlar" (without "tahlil" = not Literature Review)
                if "adabiyot" in hc and "tahlil" not in hc:
                    return True, False
                # Standalone "Литература" / "Библиография" (without "обзор")
                if ("литератур" in hc and "обзор" not in hc and "метод" not in hc) or "библиограф" in hc:
                    return True, False
                return False, False

            pending_inline_label = None  # For "Kalit so'zlar: ..." inline format

            def _is_keywords_heading(h):
                hc = re.sub(r'^[\d]+[\.\)]\s*', '', h).strip().lower()
                return hc in ("kalit so'zlar", "kalit sozlar", "keywords", "ключевые слова")

            for blk in blocks:
                # 1. Headings (explicit CENTER to override inherited JUSTIFY)
                if blk.startswith("# "):
                    h = ExportService._clean_text(blk[2:])
                    if _is_keywords_heading(h):
                        pending_inline_label = re.sub(r'^[\d]+[\.\)]\s*', '', h).strip()
                        prev_heading = pending_inline_label.lower()
                        continue
                    p = doc.add_paragraph(h)
                    p.style = doc.styles["Heading 1"]
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    prev_heading = re.sub(r'^[\d]+[\.\)]\s*', '', h).strip().lower()
                    in_references, is_combined = _check_refs_heading(h)
                    refs_started = not is_combined  # Pure refs → start immediately; combined → wait
                    if in_references:
                        ref_counter = 0
                elif blk.startswith("## "):
                    h = ExportService._clean_text(blk[3:])
                    if _is_keywords_heading(h):
                        pending_inline_label = re.sub(r'^[\d]+[\.\)]\s*', '', h).strip()
                        prev_heading = pending_inline_label.lower()
                        continue
                    p = doc.add_paragraph(h)
                    p.style = doc.styles["Heading 2"]
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    prev_heading = re.sub(r'^[\d]+[\.\)]\s*', '', h).strip().lower()
                    in_references, is_combined = _check_refs_heading(h)
                    refs_started = not is_combined
                    if in_references:
                        ref_counter = 0
                elif blk.startswith("### "):
                    h = ExportService._clean_text(blk[4:])
                    p = doc.add_paragraph(h)
                    p.style = doc.styles["Heading 3"]
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    prev_heading = re.sub(r'^[\d]+[\.\)]\s*', '', h).strip().lower()
                # 2. Table placeholder
                elif blk in _table_store:
                    ExportService._render_table(doc, _table_store[blk])
                # 3. Markdown Tables (inline detection)
                elif ExportService._is_markdown_table(blk):
                    ExportService._render_table(doc, blk)
                else:
                    # Strip duplicate heading from paragraph start
                    text = blk
                    if prev_heading:
                        lines = text.split("\n", 1)
                        first = re.sub(r'\*+', '', lines[0]).strip()
                        first_clean = re.sub(r'^[\d]+[\.\)]\s*', '', first).strip().lower()
                        if first_clean == prev_heading:
                            text = lines[1].strip() if len(lines) > 1 else ""
                    if not text.strip():
                        continue
                    # Inline keywords: "Kalit so'zlar: ..."
                    if pending_inline_label:
                        label = pending_inline_label
                        pending_inline_label = None
                        p = doc.add_paragraph()
                        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                        run = p.add_run(f"{label}: ")
                        run.bold = True
                        p.add_run(ExportService._clean_text(text))
                        continue
                    # References: numbered list
                    if in_references:
                        for line in text.strip().split("\n"):
                            line = line.strip()
                            if not line:
                                continue
                            cl = re.sub(r'^[\d]+[\.\)]\s*', '', line)
                            cl = re.sub(r'^[-•●]\s*', '', cl).strip()
                            if not cl:
                                continue
                            cl_low = cl.lower()
                            # Combined section: detect where references start
                            if not refs_started:
                                if any(x in cl_low for x in ["adabiyot", "reference", "литератур", "manba"]):
                                    refs_started = True
                                    # Skip sub-header line like "Adabiyotlar:" itself
                                    if cl_low.endswith(":") or len(cl) < 30:
                                        continue
                                elif re.search(r'\b(19|20)\d{2}\b', cl):
                                    # Line contains a year (2000, 2012...) — it's a reference entry
                                    refs_started = True
                                    # Don't continue — number this line below
                                else:
                                    # Keywords part — write as normal paragraph
                                    p = doc.add_paragraph()
                                    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                                    ExportService._apply_formatting(p, cl)
                                    continue
                            ref_counter += 1
                            if ref_counter > 10:
                                break  # MAXIMUM 10 references
                            p = doc.add_paragraph()
                            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                            ExportService._apply_formatting(p, f"{ref_counter}. {cl}")
                    # 3. Table inside else block (pipe-lines not caught by top-level check)
                    elif ExportService._is_markdown_table(text):
                        ExportService._render_table(doc, text)
                    # 4. Bullets
                    elif ExportService._looks_like_list(text):
                        for item in ExportService._split_bullets(text):
                            p = doc.add_paragraph(style="List Bullet")
                            ExportService._apply_formatting(p, item)
                    # 4. Normal Paragraph
                    else:
                        p = doc.add_paragraph()
                        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                        ExportService._apply_formatting(p, text)

            # Plagiarism Seal (Removed as per user request)
            # ExportService._add_plagiarism_seal(doc)

            doc.save(str(path))
            return True

        except Exception as e:
            logger.error(f"Docx Export Error: {e}")
            return False

    @staticmethod
    def _add_page_numbers(doc: docx.Document) -> None:
        # Faqat raqam (1, 2, 3...) — "Page" so'zisiz
        for section in doc.sections:
            footer = section.footer
            footer.is_linked_to_previous = False
            p = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER

            fld = OxmlElement("w:fldSimple")
            fld.set(qn("w:instr"), "PAGE")
            run_elem = OxmlElement("w:r")
            rpr = OxmlElement("w:rPr")
            rfont = OxmlElement("w:rFonts")
            rfont.set(qn("w:ascii"), "Times New Roman")
            rfont.set(qn("w:hAnsi"), "Times New Roman")
            rpr.append(rfont)
            sz = OxmlElement("w:sz")
            sz.set(qn("w:val"), "20")  # 10pt
            rpr.append(sz)
            run_elem.append(rpr)
            fld.append(run_elem)
            p._p.append(fld)

    # ---------------------------
    # PPTX
    # ---------------------------
    @staticmethod
    async def to_pptx(title: str, content: str, path: Path, meta: dict = None, slide_images: dict = None) -> bool:
        return await asyncio.to_thread(ExportService._to_pptx_sync, title, content, path, meta, slide_images)

    @staticmethod
    def _to_pptx_sync(title: str, content: str, path: Path, meta: dict = None, slide_images: dict = None) -> bool:
        from io import BytesIO
        from pptx.util import Inches, Emu

        try:
            prs = Presentation()
            meta = meta or {}
            slide_images = slide_images or {}

            style_mode = meta.get("ppt_style", "akademik")
            template_name = meta.get("ppt_template", "")
            clean_title = title.strip()

            # Resolve template image path
            template_path = None
            if template_name:
                from pathlib import Path as _P
                _templates_dir = _P(__file__).parent.parent / "templates" / "presentations"
                _candidate = _templates_dir / f"template_{template_name}.png"
                # Try numbered variants
                if not _candidate.exists():
                    for _f in _templates_dir.glob(f"template_*_{template_name}.png"):
                        _candidate = _f
                        break
                if _candidate.exists():
                    template_path = str(_candidate)

            def _apply_bg(slide):
                if not template_path:
                    return
                try:
                    pic = slide.shapes.add_picture(template_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
                    # Move background picture to bottom (behind all shapes)
                    spTree = pic._element.getparent()
                    spTree.remove(pic._element)
                    spTree.insert(2, pic._element)
                except Exception:
                    pass

            # Title slide
            slide0 = prs.slides.add_slide(prs.slide_layouts[0])
            _apply_bg(slide0)
            slide0.shapes.title.text = clean_title
            slide0.placeholders[1].text = f"Tayyorlandi: Akademik Bot\nUslub: {style_mode.capitalize()}"
            # Make title text white for dark templates
            if template_path:
                from pptx.dml.color import RGBColor as _PptxRGB
                for ph in (slide0.shapes.title, slide0.placeholders[1]):
                    for para in ph.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.color.rgb = _PptxRGB(0xFF, 0xFF, 0xFF)

            font_name = "Times New Roman"
            if style_mode == "biznes":
                font_name = "Arial"
            elif style_mode == "kreativ":
                font_name = "Verdana"

            slides = ExportService._parse_slides(content or "")
            if not slides:
                s = prs.slides.add_slide(prs.slide_layouts[1])
                s.shapes.title.text = "Mazmun"
                s.placeholders[1].text = (content or "")[:1000]
                prs.save(str(path))
                return True

            # Slayd o'lchamlari
            slide_width = prs.slide_width   # 9144000 EMU = 10 inches
            slide_height = prs.slide_height  # 6858000 EMU = 7.5 inches

            # Filter out duplicate title slides from AI (slide_0 already has title)
            filtered_slides = []
            for idx, (s_title, bullets, notes) in enumerate(slides):
                _t = re.sub(r'^\s*\d+[\.\)]\s*', '', ExportService._clean_text(s_title)).strip().lower()
                # Skip AI's first title slide since prs has one already
                if idx == 0 and any(k in _t for k in ["titul slayd", "title slide", "titul"]):
                    continue
                filtered_slides.append((s_title, bullets, notes))

            for idx, (s_title, bullets, notes) in enumerate(filtered_slides):
                s = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                _apply_bg(s)
                from pptx.util import Inches
                from pptx.dml.color import RGBColor as PptxRGB

                has_image = False  # AI images disabled
                title_text = ExportService._clean_text(s_title)
                # Remove leading numbering like "1.", "10.", "2)"
                title_text = re.sub(r'^\s*\d+[\.\)]\s*', '', title_text).strip()
                # White text on dark template, dark text otherwise
                text_color = PptxRGB(0xFF, 0xFF, 0xFF) if template_path else PptxRGB(30, 30, 30)

                # Detect thanks slides
                _low = title_text.lower()
                is_thanks = any(k in _low for k in ["rahmat", "thank", "спасибо", "e'tiboringiz"])

                if is_thanks:
                    # Center big text, no bullets
                    big_box = s.shapes.add_textbox(
                        left=Inches(0.5), top=Inches(2.8),
                        width=Inches(9.0), height=Inches(2.0)
                    )
                    _tf = big_box.text_frame
                    _tf.word_wrap = True
                    _p = _tf.paragraphs[0]
                    _p.alignment = 2  # CENTER
                    _p.text = title_text
                    _p.font.size = PptxPt(40)
                    _p.font.name = font_name
                    _p.font.bold = True
                    _p.font.color.rgb = text_color
                    continue

                # --- TITLE (yuqori burchakda, chiziq tepasida) ---
                title_shape = s.shapes.add_textbox(
                    left=Inches(0.7), top=Inches(0.15),
                    width=Inches(8.6), height=Inches(0.75)
                )
                tf = title_shape.text_frame
                tf.word_wrap = True
                tp = tf.paragraphs[0]
                tp.text = title_text
                tp.font.size = PptxPt(28)
                tp.font.name = font_name
                tp.font.bold = True
                tp.font.color.rgb = text_color

                # --- BULLETS (to'liq kenglik, katta shrift) ---
                bullet_box = s.shapes.add_textbox(
                    left=Inches(0.9), top=Inches(1.4),
                    width=Inches(8.2), height=Inches(5.6)
                )
                body = bullet_box.text_frame
                body.word_wrap = True

                clean_bullets = []
                for b in bullets[:5]:
                    txt = ExportService._clean_text(b).strip()
                    # Remove leading *, **, -, • that AI sometimes adds
                    txt = re.sub(r'^[\*\-•‣◦\s]+', '', txt).strip()
                    # Remove stray ** markers inside
                    txt = txt.replace("**", "").replace("__", "")
                    if not txt:
                        continue
                    if len(txt) > 180:
                        txt = txt[:177] + "..."
                    clean_bullets.append(txt)

                # Dinamik shrift — bullet soniga qarab
                n = len(clean_bullets)
                if n <= 3:
                    bullet_size, spacing = 22, 16
                elif n == 4:
                    bullet_size, spacing = 20, 12
                else:
                    bullet_size, spacing = 18, 10

                for i, b in enumerate(clean_bullets):
                    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                    p.text = f"\u2022  {b}"
                    p.font.size = PptxPt(bullet_size)
                    p.font.name = font_name
                    p.font.color.rgb = text_color
                    p.space_after = PptxPt(spacing)

                # --- RASM (o'ng tomonda, title ostida) ---
                if has_image:
                    try:
                        img_stream = BytesIO(slide_images[idx])
                        pic = s.shapes.add_picture(
                            img_stream,
                            left=Inches(5.7), top=Inches(1.5),
                            width=Inches(4.0),
                        )
                        # Rasm border (yupqa kulrang chiziq)
                        pic.line.color.rgb = PptxRGB(200, 200, 200)
                        pic.line.width = PptxPt(1)
                    except Exception as e:
                        logger.warning(f"Slide {idx} image failed: {e}")

                # Speaker notes
                if notes:
                    note_tf = s.notes_slide.notes_text_frame
                    note_tf.text = ExportService._clean_text(notes)

            prs.save(str(path))
            return True

        except Exception as e:
            logger.error(f"Pptx Export Error: {e}")
            return False

    # ---------------------------
    # Parsing helpers
    # ---------------------------
    @staticmethod
    # Structured Metadata is now passed via 'meta' argument, these parsers are obsolete.

    @staticmethod
    def _parse_slides(content: str):
        """
        Expects:
        --- SLAYD X: Title ---
        Bullet...
        Speaker Notes: ....
        """
        pattern = r"---\s*SLAYD\s*\d+\s*:?\s*(.*?)\s*---"
        parts = re.split(pattern, content)

        if len(parts) <= 1:
            return []

        slides = []
        for i in range(1, len(parts), 2):
            title = parts[i].strip()
            body = parts[i + 1].strip() if i + 1 < len(parts) else ""

            bullets, notes = ExportService._split_body_notes(body)
            slides.append((title, bullets, notes))

        return slides

    @staticmethod
    def _split_body_notes(text: str):
        # Try separate Speaker Notes
        notes = ""
        m = re.search(r"(speaker\s*notes\s*:)(.*)$", text, flags=re.I | re.S)
        if m:
            notes = m.group(2).strip()
            text = text[: m.start()].strip()

        # Strip AI meta-labels like "*Slayd sarlavhasi: ...", "*Slayd mazmuni:"
        cleaned_lines = []
        for ln in text.splitlines():
            stripped = ln.strip()
            # Skip meta-label lines
            if re.match(r"^\*?\s*slayd\s+(sarlavhasi|mazmuni)\s*:", stripped, re.I):
                continue
            cleaned_lines.append(ln)
        text = "\n".join(cleaned_lines)

        bullets = ExportService._split_bullets(text)
        if not bullets:
            bullets = [text.strip()] if text.strip() else []
        return bullets, notes

    @staticmethod
    def _looks_like_list(text: str) -> bool:
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
        if len(lines) < 2:
            return False
        hits = 0
        for ln in lines:
            if ln.startswith(("-", "•", "*")) or re.match(r"^\d+[\.\)]\s+", ln):
                hits += 1
        return hits >= 2

    @staticmethod
    def _is_markdown_table(blk: str) -> bool:
        lines = [ln.strip() for ln in blk.strip().splitlines() if ln.strip()]
        if len(lines) < 2: return False
        # Must have | in first line
        if "|" not in lines[0]: return False
        # Check if any line is a separator like |---|---|
        for ln in lines:
            if re.match(r"^[ \t]*\|?[\-\s\|:]+\|?[ \t]*$", ln):
                return True
        # Or: if majority of lines contain |, treat as table
        pipe_lines = sum(1 for ln in lines if "|" in ln)
        if pipe_lines >= 2 and pipe_lines >= len(lines) * 0.6:
            return True
        return False

    @staticmethod
    def _render_table(doc, blk: str):
        lines = [ln.strip() for ln in blk.strip().splitlines() if ln.strip()]
        table_data = []
        caption_lines = []

        for ln in lines:
            # Skip separator line
            if re.match(r"^[ \t]*\|?[\-\s\|:]+\|?[ \t]*$", ln):
                continue
            # Lines without | are table captions/titles
            if "|" not in ln:
                caption_lines.append(ln)
                continue
            # Remove leading/trailing pipes and split
            cells = [c.strip() for c in ln.strip("|").split("|")]
            if cells:
                table_data.append(cells)

        # Render caption above table
        for cap in caption_lines:
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run(ExportService._clean_text(cap))
            run.bold = True
        
        if not table_data:
            return

        num_rows = len(table_data)
        num_cols = max(len(r) for r in table_data)

        table = doc.add_table(rows=num_rows, cols=num_cols)
        table.style = 'Table Grid'

        # Auto-fit column widths based on page width
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        page_width_cm = 16.0  # A4 usable width
        col_width = Cm(page_width_cm / num_cols)

        for col in table.columns:
            for cell in col.cells:
                cell.width = col_width

        for r_idx, row_cells in enumerate(table_data):
            for c_idx, cell_value in enumerate(row_cells):
                if c_idx < num_cols:
                    cell = table.cell(r_idx, c_idx)
                    cell.width = col_width
                    # Light cleaning — preserve short values and numeric data
                    cv = (cell_value or "").strip()
                    cv = cv.replace("**", "").replace("__", "").replace("*", "")
                    p = cell.paragraphs[0]
                    p.clear()
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    run = p.add_run(cv)
                    run.font.size = DocxPt(10)
                    # Bold headers
                    if r_idx == 0:
                        run.bold = True

    @staticmethod
    def _split_bullets(text: str) -> list[str]:
        lines = [ln.strip() for ln in (text or "").splitlines() if ln.strip()]
        out = []
        for ln in lines:
            ln2 = re.sub(r"^(\-|\*|•)\s*", "", ln).strip()
            ln2 = re.sub(r"^\d+[\.\)]\s*", "", ln2).strip()
            if ln2:
                out.append(ln2)
        return out

    # ---------------------------
    # Title page (your logic kept)
    # ---------------------------
    @staticmethod
    def _render_title_page(doc, title, meta, doc_type: str = "coursework"):
        lang = (meta.get("language") or "uz").lower()

        # Language-specific labels
        _labels = {
            "uz": {
                "ministry": "O'ZBEKISTON RESPUBLIKASI OLIY TA'LIM, FAN VA INNOVATSIYALAR VAZIRLIGI",
                "uni_default": "UNIVERSITET NOMI",
                "types": {
                    "article": "ILMIY MAQOLA", "coursework": "KURS ISHI",
                    "independent": "MUSTAQIL ISH", "diploma": "BITIRUV MALAKAVIY ISHI",
                    "dissertation": "DISSERTATSIYA", "manual": "O'QUV QO'LLANMA",
                },
                "type_default": "KURS ISHI",
                "subject_suffix": "fanidan",
                "topic_label": "MAVZU",
                "completed_by": "Bajardi:",
                "student": "Talaba",
                "advisor": "Ilmiy rahbar:",
                "city": "Toshkent",
            },
            "en": {
                "ministry": "MINISTRY OF HIGHER EDUCATION, SCIENCE AND INNOVATION OF THE REPUBLIC OF UZBEKISTAN",
                "uni_default": "UNIVERSITY NAME",
                "types": {
                    "article": "SCIENTIFIC ARTICLE", "coursework": "COURSEWORK",
                    "independent": "INDEPENDENT WORK", "diploma": "GRADUATION THESIS",
                    "dissertation": "DISSERTATION", "manual": "TEXTBOOK",
                },
                "type_default": "COURSEWORK",
                "subject_suffix": "subject",
                "topic_label": "TOPIC",
                "completed_by": "Completed by:",
                "student": "Student",
                "advisor": "Scientific advisor:",
                "city": "Tashkent",
            },
            "ru": {
                "ministry": "МИНИСТЕРСТВО ВЫСШЕГО ОБРАЗОВАНИЯ, НАУКИ И ИННОВАЦИЙ РЕСПУБЛИКИ УЗБЕКИСТАН",
                "uni_default": "НАЗВАНИЕ УНИВЕРСИТЕТА",
                "types": {
                    "article": "НАУЧНАЯ СТАТЬЯ", "coursework": "КУРСОВАЯ РАБОТА",
                    "independent": "САМОСТОЯТЕЛЬНАЯ РАБОТА", "diploma": "ВЫПУСКНАЯ КВАЛИФИКАЦИОННАЯ РАБОТА",
                    "dissertation": "ДИССЕРТАЦИЯ", "manual": "УЧЕБНОЕ ПОСОБИЕ",
                },
                "type_default": "КУРСОВАЯ РАБОТА",
                "subject_suffix": "по предмету",
                "topic_label": "ТЕМА",
                "completed_by": "Выполнил:",
                "student": "Студент",
                "advisor": "Научный руководитель:",
                "city": "Ташкент",
            },
        }
        L = _labels.get(lang, _labels["uz"])

        # Ministry Header
        p = doc.add_paragraph(L["ministry"])
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.runs[0].bold = True
        p.runs[0].font.size = DocxPt(12)

        otm = meta.get("uni", L["uni_default"]).upper()
        p = doc.add_paragraph(otm)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.runs[0].bold = True
        p.paragraph_format.space_after = DocxPt(24)

        for _ in range(2):
            doc.add_paragraph("")

        # Type Indicator
        type_str = L["types"].get(doc_type, L["type_default"])

        p = doc.add_paragraph(type_str)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.runs[0].bold = True
        p.runs[0].font.size = DocxPt(20)
        p.paragraph_format.space_after = DocxPt(12)

        subj = meta.get("subject", "").upper()
        if subj:
            p = doc.add_paragraph(f"\"{subj}\" {L['subject_suffix']}")
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p.runs[0].font.size = DocxPt(14)

        for _ in range(2):
            doc.add_paragraph("")

        p = doc.add_paragraph(f"{L['topic_label']}: {title}")
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.runs[0].bold = True
        p.runs[0].font.size = DocxPt(16)

        for _ in range(6):
            doc.add_paragraph("")

        # Student/Advisor Table
        student_name = meta.get("student_name") or L["student"]
        advisor = meta.get("advisor", "................")

        # Spacer for table
        doc.add_paragraph("").paragraph_format.first_line_indent = Cm(0)

        table = doc.add_table(rows=2, cols=2)
        table.alignment = WD_ALIGN_PARAGRAPH.RIGHT

        table.cell(0, 0).text = L["completed_by"]
        table.cell(0, 1).text = student_name
        table.cell(1, 0).text = L["advisor"]
        table.cell(1, 1).text = advisor

        # Right align all cells in the right table
        for row in table.rows:
            for cell in row.cells:
                cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.LEFT

        for _ in range(4):
            doc.add_paragraph("")

        p = doc.add_paragraph(f"{L['city']} — {datetime.now().year}")
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    @staticmethod
    def _render_article_header(doc, title: str, meta: dict):
        """Maqola — muallif ma'lumotlari (o'ng tomonda) + sarlavha."""
        # Muallif ma'lumotlari — o'ng tomonda
        authors = meta.get("authors", "")
        workplace = meta.get("workplace", "")
        author_email = meta.get("author_email", "")
        advisor = meta.get("advisor", "")

        info_lines = []
        if authors:
            info_lines.append(authors)
        if workplace:
            info_lines.append(workplace)
        if author_email:
            info_lines.append(author_email)
        if advisor:
            info_lines.append(f"Ilmiy rahbar: {advisor}")

        # Sarlavha (bold, centered) — BIRINCHI
        p = doc.add_paragraph(ExportService._clean_text(title).upper())
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.first_line_indent = Cm(0)
        p.paragraph_format.space_after = DocxPt(12)
        run = p.runs[0]
        run.bold = True
        run.font.size = DocxPt(14)

        # Muallif ma'lumotlari — o'rtada (sarlavhadan keyin)
        if info_lines:
            for line in info_lines:
                p = doc.add_paragraph(line)
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                p.paragraph_format.first_line_indent = Cm(0)
                p.paragraph_format.space_after = DocxPt(0)
                p.paragraph_format.space_before = DocxPt(0)
                run = p.runs[0]
                run.italic = True
                run.font.size = DocxPt(12)
            # Bo'sh qator
            doc.add_paragraph("").paragraph_format.space_after = DocxPt(6)

    @staticmethod
    def _render_thesis_header(doc, title: str, meta: dict):
        """Tezis (konferensiya tezisi) — compact inline header, no title page."""
        udc = meta.get("udc", "")
        if udc:
            p = doc.add_paragraph(f"UO'K {udc}")
            p.alignment = WD_ALIGN_PARAGRAPH.LEFT
            p.runs[0].font.size = DocxPt(12)
            p.paragraph_format.space_after = DocxPt(6)
            p.paragraph_format.first_line_indent = Cm(0)

        # Muallif ma'lumotlari — o'ng tomonda
        authors = meta.get("authors", "")
        workplace = meta.get("workplace", "")
        author_email = meta.get("author_email", "")
        advisor = meta.get("advisor", "")

        info_lines = []
        if authors:
            info_lines.append(authors)
        if workplace:
            info_lines.append(workplace)
        if author_email:
            info_lines.append(author_email)
        if advisor:
            info_lines.append(f"Ilmiy rahbar: {advisor}")

        if info_lines:
            for line in info_lines:
                p = doc.add_paragraph(line)
                p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                p.paragraph_format.first_line_indent = Cm(0)
                p.paragraph_format.space_after = DocxPt(0)
                p.paragraph_format.space_before = DocxPt(0)
                run = p.runs[0]
                run.italic = True
                run.font.size = DocxPt(12)
            doc.add_paragraph("").paragraph_format.space_after = DocxPt(6)

        # Title (bold, centered)
        p = doc.add_paragraph(ExportService._clean_text(title).upper())
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.first_line_indent = Cm(0)
        p.paragraph_format.space_after = DocxPt(12)
        run = p.runs[0]
        run.bold = True
        run.font.size = DocxPt(14)

    # ---------------------------
    # Formatting helpers
    # ---------------------------
    @staticmethod
    def _add_plagiarism_seal(doc):
        doc.add_page_break()
        
        # Seal Header
        p = doc.add_paragraph("TADQIQOTNING NOYOB LIK (UNIQUENESS) HISOBOTI")
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.runs[0].bold = True
        p.runs[0].font.size = DocxPt(16)
        
        for _ in range(2):
            doc.add_paragraph("")
            
        import random
        score = random.randint(92, 98)
        
        table = doc.add_table(rows=4, cols=2)
        table.style = 'Table Grid'
        
        table.cell(0, 0).text = "Tekshiruv tizimi:"
        table.cell(0, 1).text = "Stealth AI & Plagiarism Guard v4.2"
        
        table.cell(1, 0).text = "Noyoblik darajasi:"
        table.cell(1, 1).text = f"{score}% (YUQORI NOYOBLIK)"
        
        table.cell(2, 0).text = "Sifat nazorati (Critique):"
        table.cell(2, 1).text = "Muvaffaqiyatli o'tdi"
        
        table.cell(3, 0).text = "Sana:"
        table.cell(3, 1).text = datetime.now().strftime("%d.%m.%Y")

        for _ in range(2):
            doc.add_paragraph("")

        p = doc.add_paragraph("Eslatma: Ushbu hisobot avtomatik ravishda shakllantirildi va mazkur tadqiqot AI chattlardan butunlay tozalanganligini bildiradi.")
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.runs[0].italic = True
        p.runs[0].font.size = DocxPt(10)

    @staticmethod
    def _apply_formatting(paragraph, text: str):
        """
        Parses basic markdown (**bold**, *italic*) and applies it to the docx paragraph.
        """
        if not text:
            return

        # Basic **bold** and *italic* parser
        # Clean common markdown artifacts but keep valid formatting markers for the regex
        text = text.replace(":*", ": ").replace("*.", ". ").replace("*,", ", ")
        # Remove orphan trailing/leading asterisks (e.g. "Natijalar*" → "Natijalar")
        text = re.sub(r'(\w)\*(\s|$)', r'\1\2', text)  # trailing *
        text = re.sub(r'(^|\s)\*(\w)', r'\1\2', text)   # leading * before word (not bold)
        
        pattern = r"(\*\*.*?\*\*|\*.*?\*)"
        parts = re.split(pattern, text or "")
        for part in parts:
            if part.startswith("**") and part.endswith("**"):
                # Double clean internal stars
                clean_part = part[2:-2].replace("*", "").strip()
                paragraph.add_run(clean_part).bold = True
            elif part.startswith("*") and part.endswith("*") and len(part) > 2:
                clean_part = part[1:-1].replace("*", "").strip()
                paragraph.add_run(clean_part).italic = True
            else:
                paragraph.add_run(part)

    @staticmethod
    def _clean_text(text: str) -> str:
        """Enhanced text cleaning with aggressive AI artifact removal."""
        if not text:
            return ""
        
        # 0. Detect "technical" context (formulas, code, keywords)
        tech_indicators = [
            r"python", r"sql", r"code", r"algoritm", r"formula", 
            r"tenglama", r"da'stur", r"matematik", r"fizik", r"kimyo",
            r"\[.*?\]", r"\{.*?\}", r"\(\d+\)"
        ]
        is_technical = any(re.search(ind, text.lower()) for ind in tech_indicators)
        
        # 1. Strip possible markdown code fences
        text = re.sub(r"```[a-zA-Z]*\n?", "", text)
        text = text.replace("```", "")
        
        # 2. Character cleaning - Remove markdown artifacts
        banned_symbols = [
             r'###', r'##', r'\*\*', r'__', r'==', r'~~', r'«', r'»', r'…', 
             r'\?\?\?', r'\*\*\*', r'//', r'%%', r'`'
        ]
        
        # Only strip curly/square brackets if NOT technical
        if not is_technical:
            banned_symbols.extend([r'\{', r'\}', r'\[', r'\]'])
            
        for symbol in banned_symbols:
            text = re.sub(symbol, "", text)
            
        # Remove leftover markdown headers specifically
        text = re.sub(r"^#+\s*", "", text, flags=re.MULTILINE)
        
        # 3. List markers at start of lines (- , * , 1. ) but keep text
        text = re.sub(r"^[ \t]*[\-\*\u2022]\s+", "", text, flags=re.MULTILINE)
        
        # 4. Cleanup escaped markdown artifacts
        text = text.replace("\\-", "-").replace("\\*", "*").replace("\\.", ".")
        
        # 5. Remove WORD COUNT REPORT and similar AI meta-text from final export
        text = re.sub(r'WORD\s+COUNT\s+REPORT:.*', '', text, flags=re.IGNORECASE | re.DOTALL)
        text = re.sub(r'Total\s+words\s*\(excluding\s+references\)\s*:\s*\d+.*', '', text, flags=re.IGNORECASE | re.DOTALL)
        text = re.sub(r'^Status\s*:\s*(?:PASS|FAIL)\s*$', '', text, flags=re.IGNORECASE | re.MULTILINE)
        text = re.sub(r'^Target\s+range\s*:\s*[\d–\-]+\s*$', '', text, flags=re.IGNORECASE | re.MULTILINE)
        
        # 6. Remove common AI leak phrases / chatter (ENHANCED)
        ai_phrases_patterns = [
            # Uzbek
            r"\b(shuningdek|shunday qilib|xulosa qilib aytganda),?\s*",
            r"\bushbu bo'limda\b",
            r"\bmana maqola\b",
            r"\bumid qilamanki\b",
            r"\bquyidagicha\b",
            r"\bbu yerda\b",
            r"\balbatta,\s*",
            r"\bavvalambor,\s*",
            # English
            r"\bhere is the\b",
            r"\bI hope this\b",
            r"\bas requested\b",
            r"\blet me\b",
            r"\bI will\b",
            r"\bplease note\b",
            # Russian
            r"\bвот статья\b",
            r"\bнадеюсь\b",
            r"\bкак запрошено\b",
            # Placeholders
            r"\blorem\b",
            r"\bplaceholder\b",
            r"\bxxx\b",
            r"\bsample\b",
            r"___+",
            r"\.\.\.+",
        ]
        
        for pattern in ai_phrases_patterns:
            text = re.sub(pattern, "", text, flags=re.IGNORECASE)
        
        # 7. Remove meta-commentary lines
        meta_commentary = [
            r"^.*?(bu bo'lim|this section|этот раздел).*?$",
            r"^.*?(maqolaning mazmuni|article content|содержание статьи).*?$",
            r"^.*?(quyida|below|ниже).*?(keltirilgan|presented|представлено).*?$",
        ]
        
        for pattern in meta_commentary:
            text = re.sub(pattern, "", text, flags=re.MULTILINE | re.IGNORECASE)
        
        # 8. Clean up lines and rebuild
        lines = []
        for line in text.splitlines():
            line_strip = line.strip()
            if not line_strip: 
                continue
            
            # Skip lines that are only punctuation or very short
            if len(line_strip) < 3:
                continue
            
            # Skip lines with only special characters (but keep numeric data like "2150.4 ± 187.2")
            if re.match(r'^[^a-zA-Zа-яА-ЯёЁўЎқҚғҒҳҲ0-9]+$', line_strip):
                continue
            
            lines.append(line_strip)
            
        return "\n".join(lines).strip()


export_service = ExportService()
